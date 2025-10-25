# utils/file_utils.py
import os
import cv2
import json
import hashlib
import imagehash
import numpy as np
import pandas as pd
from io import BytesIO
from PIL import Image
from collections import defaultdict
from numpy.linalg import norm

# Handling imports for dependencies
try:
    LANCZOS = Image.Resampling.LANCZOS
except AttributeError:
    LANCZOS = Image.LANCZOS # type: ignore

try:
    import cv2
    VIDEO_SUPPORT = True
except ImportError:
    VIDEO_SUPPORT = False
    print("Warning: opencv-python not installed. Video duplicate detection disabled.")

try:
    import librosa
    AUDIO_SUPPORT = True
except ImportError:
    AUDIO_SUPPORT = False
    print("Warning: librosa not installed. Audio duplicate detection disabled.")

try:
    from PyPDF2 import PdfReader
    from sentence_transformers import SentenceTransformer
    DOC_SUPPORT = True
    text_model = SentenceTransformer("/app/models/all-MiniLM-L6-v2")
except ImportError:
    DOC_SUPPORT = False
    text_model = None
    print("Warning: PyPDF2/sentence-transformers not installed. Document duplicate detection limited.")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

print(f"VIDEO_SUPPORT: {VIDEO_SUPPORT}")
print(f"AUDIO_SUPPORT: {AUDIO_SUPPORT}")
print(f"DOC_SUPPORT: {DOC_SUPPORT}")
print(f"PPTX_SUPPORT: {PPTX_SUPPORT}")


EXTENSIONS = {
    "images": ['.jpg', '.jpeg', '.png', '.gif', '.svg', '.webp', '.bmp'],
    "videos": ['.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.webm'],
    "audios": ['.mp3', '.wav', '.aac', '.flac', '.ogg', '.m4a', '.wma'],
    "pdfs": ['.pdf'],
    "documents": ['.doc', '.docx', '.txt'],
    "tables": ['.csv', '.xls', '.xlsx'],
    "pptx": ['.ppt', '.pptx'],
    "others": ['.exe', '.bin', '.zip', '.rar', '.7z']
}


SIMILARITY_THRESHOLD = 20 


def get_file_type(filename):
    """Determine file type from extension."""
    ext = os.path.splitext(filename.lower())[1]
    for file_type, extensions in EXTENSIONS.items():
        if ext in extensions:
            return file_type
    return "others"


def compute_exact_hash(file_bytes):
    """Compute MD5 hash for exact duplicate detection."""
    return hashlib.md5(file_bytes).hexdigest()


def serialize_hash(hash_array):
    """Convert hash arrays to JSON-serializable format."""
    if isinstance(hash_array, np.ndarray):
        return hash_array.astype(np.float64).tolist()
    elif isinstance(hash_array, (np.number, np.floating, np.integer)):
        return float(hash_array)
    elif isinstance(hash_array, list):
        return [serialize_hash(item) for item in hash_array]
    return hash_array


def hash_image(file_bytes):
    """Hash image using multiple perceptual hashes for better matching."""
    try:
        image = Image.open(BytesIO(file_bytes))
        
        phash_val = str(imagehash.phash(image))
        ahash_val = str(imagehash.average_hash(image))
        dhash_val = str(imagehash.dhash(image))
        
        image_rgb = image.convert("RGB")
        img_array = np.array(image_rgb)
        
        hist_r = np.histogram(img_array[:,:,0], bins=8, range=(0,256))[0]
        hist_g = np.histogram(img_array[:,:,1], bins=8, range=(0,256))[0]
        hist_b = np.histogram(img_array[:,:,2], bins=8, range=(0,256))[0]
        color_hist = np.concatenate([hist_r, hist_g, hist_b])
        color_hist = color_hist / (color_hist.sum() + 1e-10)  
        
        combined_hash = {
            'phash': phash_val,
            'ahash': ahash_val,
            'dhash': dhash_val,
            'color_hist': serialize_hash(color_hist)
        }
        
        return json.dumps(combined_hash)
    except Exception as e:
        print(f"Error hashing image: {e}")
        return None


def hash_video(file_bytes, num_frames=10, resize_dim=(256, 256)):
    """Hash video using ORB descriptors from sampled frames."""
    if not VIDEO_SUPPORT:
        return None
    
    import tempfile
    import os as os_module
    
    temp_file = None
    cap = None
    
    try:
        fd, temp_path = tempfile.mkstemp(suffix='.mp4')
        
        try:
            os_module.write(fd, file_bytes)
            os_module.close(fd)  
            
            cap = cv2.VideoCapture(temp_path)
            
            if not cap.isOpened():
                print("Error: Could not open video file")
                return None
            
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            
            if total_frames <= 0:
                print(f"⚠️  Video has {total_frames} frames, trying sequential read...")
                
                frame_count = 0
                max_attempts = num_frames * 10
                all_descriptors = []
                orb = cv2.ORB_create() # type: ignore
                
                while frame_count < max_attempts:
                    success, frame = cap.read()
                    if not success:
                        break
                    
                    if frame_count % max(1, max_attempts // num_frames) == 0:
                        try:
                            resized = cv2.resize(frame, resize_dim)
                            gray = cv2.cvtColor(resized, cv2.COLOR_BGR2GRAY)
                            keypoints, descriptors = orb.detectAndCompute(gray, None)
                            if descriptors is not None and len(descriptors) > 0:
                                all_descriptors.append(descriptors)
                        except Exception as e:
                            print(f"Error processing frame {frame_count}: {e}")
                    
                    frame_count += 1
                
                cap.release()
                
                if all_descriptors:
                    all_descriptors = np.vstack(all_descriptors)
                    descriptor_hash = np.mean(all_descriptors, axis=0)
                    print(f"   ✅ Processed {len(all_descriptors)} frames via sequential read")
                    return serialize_hash(descriptor_hash)
                else:
                    print(f"   ⚠️  No valid frames extracted")
                    return None
            
            print(f"   📹 Video info: {total_frames} frames @ {fps:.2f} fps")
            
            orb = cv2.ORB_create() # type: ignore
            all_descriptors = []
            
            frame_indices = np.linspace(0, total_frames - 1, num_frames, dtype=np.int32)
            frames_processed = 0
            
            for idx in frame_indices:
                try:
                    cap.set(cv2.CAP_PROP_POS_FRAMES, int(idx))
                    success, frame = cap.read()
                    
                    if success and frame is not None:
                        resized = cv2.resize(frame, resize_dim)
                        gray = cv2.cvtColor(resized, cv2.COLOR_BGR2GRAY)
                        keypoints, descriptors = orb.detectAndCompute(gray, None)
                        if descriptors is not None and len(descriptors) > 0:
                            all_descriptors.append(descriptors)
                            frames_processed += 1
                except Exception as e:
                    print(f"Error processing frame at index {idx}: {e}")
                    continue
            
            cap.release()
            
            print(f"   ✅ Processed {frames_processed}/{len(frame_indices)} frames successfully")
            
            if all_descriptors:
                all_descriptors = np.vstack(all_descriptors)
                descriptor_hash = np.mean(all_descriptors, axis=0)
                return serialize_hash(descriptor_hash)
            else:
                print(f"   ⚠️  No valid descriptors extracted")
                return None
        
        finally:
            try:
                os_module.unlink(temp_path)
            except Exception as e:
                print(f"Error deleting temp file: {e}")
        
    except Exception as e:
        print(f"Error hashing video: {e}")
        import traceback
        traceback.print_exc()
        return None
    finally:
        if cap is not None:
            cap.release()


def hash_audio(file_bytes):
    """Hash audio using MFCC features."""
    if not AUDIO_SUPPORT:
        return None
    
    try:
        y, sr = librosa.load(BytesIO(file_bytes), sr=None, mono=True, duration=60)
        
        if len(y) == 0:
            print("Error: Audio file is empty")
            return None
        
        mfcc = librosa.feature.mfcc(y=y, sr=sr, n_mfcc=13)
        mfcc_mean = np.mean(mfcc, axis=1)
        
        return serialize_hash(mfcc_mean)
    except Exception as e:
        print(f"Error hashing audio: {e}")
        return None


def hash_text_file(file_bytes):
    """Hash text file using sentence embeddings."""
    if not DOC_SUPPORT or text_model is None:
        return None
    
    try:
        content = file_bytes.decode('utf-8', errors='ignore')
        
        if not content.strip():
            print("Error: Text file is empty")
            return None
        
        embedding = text_model.encode(content, convert_to_numpy=True)
        return serialize_hash(embedding)
    except Exception as e:
        print(f"Error hashing text: {e}")
        return None


def hash_pdf_file(file_bytes):
    """Hash PDF using sentence embeddings of extracted text."""
    if not DOC_SUPPORT or text_model is None:
        return None
    
    try:
        reader = PdfReader(BytesIO(file_bytes))
        text = ""
        
        for page in reader.pages[:10]:
            try:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + " "
            except Exception as e:
                print(f"Error extracting page text: {e}")
                continue
        
        if not text.strip():
            print("Error: No text extracted from PDF")
            return None
        
        embedding = text_model.encode(text, convert_to_numpy=True)
        return serialize_hash(embedding)
    except Exception as e:
        print(f"Error hashing PDF: {e}")
        return None


def hash_table_file(file_bytes, filename):
    """Hash table files (CSV, Excel) using sentence embeddings."""
    if not DOC_SUPPORT or text_model is None:
        return None
    
    try:
        ext = os.path.splitext(filename.lower())[1]
        
        if ext == ".csv":
            df = pd.read_csv(BytesIO(file_bytes), dtype=str, encoding='utf-8', errors='ignore') # type: ignore
        else: 
            df = pd.read_excel(BytesIO(file_bytes), dtype=str, engine='openpyxl')
        
        if df.empty:
            print("Error: Table file is empty")
            return None
        
        df = df.fillna('')
        text = ' '.join(df.astype(str).values.flatten()[:10000])  
        
        if not text.strip():
            print("Error: No text content in table")
            return None
        
        embedding = text_model.encode(text, convert_to_numpy=True)
        return serialize_hash(embedding)
    except Exception as e:
        print(f"Error hashing table: {e}")
        return None


def extract_text_from_pptx(file_bytes):
    """Extract text from PPTX file."""
    if not PPTX_SUPPORT:
        return None
    
    try:
        ppt = Presentation(BytesIO(file_bytes))
        all_text = ""
        
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    try:
                        for paragraph in shape.text_frame.paragraphs: # type: ignore
                            for run in paragraph.runs:
                                all_text += run.text + " "
                    except Exception as e:
                        print(f"Error extracting text from shape: {e}")
                        continue
        
        return all_text.strip() if all_text.strip() else None
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return None


def extract_images_from_pptx(file_bytes):
    """Extract images from PPTX file."""
    if not PPTX_SUPPORT:
        return []
    
    try:
        ppt = Presentation(BytesIO(file_bytes))
        images = []
        
        for slide in ppt.slides:
            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_stream = shape.image.blob # type: ignore
                        image = Image.open(BytesIO(image_stream))
                        images.append(image)
                except Exception as e:
                    print(f"Error extracting image from slide: {e}")
                    continue
        
        return images
    except Exception as e:
        print(f"Error extracting PPTX images: {e}")
        return []


def hash_pptx_file(file_bytes):
    """Hash PPTX using text embeddings and image hashes."""
    if not PPTX_SUPPORT or not DOC_SUPPORT or text_model is None:
        return None
    
    try:
        text_content = extract_text_from_pptx(file_bytes)
        if text_content:
            text_embedding = text_model.encode(text_content, convert_to_numpy=True)
            text_embedding = serialize_hash(text_embedding)
        else:
            text_embedding = []
        
        images = extract_images_from_pptx(file_bytes)
        image_hashes = []
        for image in images:
            try:
                hash_val = str(imagehash.phash(image))
                image_hashes.append(hash_val)
            except Exception as e:
                print(f"Error hashing PPTX image: {e}")
                continue
        
        combined_hash = []
        if isinstance(text_embedding, list):
            combined_hash.extend(text_embedding)
        combined_hash.extend(image_hashes)
        
        return combined_hash if combined_hash else None
    except Exception as e:
        print(f"Error hashing PPTX: {e}")
        return None


def compute_file_hash(file_bytes, filename):
    """Compute appropriate hash based on file type."""
    ext = os.path.splitext(filename.lower())[1].strip('.')
    file_type = get_file_type(filename)
    
    hash_result = None
    
    if file_type == "images":
        hash_result = hash_image(file_bytes)
    
    elif file_type == "videos":
        hash_result = hash_video(file_bytes)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    elif file_type == "audios":
        hash_result = hash_audio(file_bytes)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    elif file_type == "documents":
        hash_result = hash_text_file(file_bytes)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    elif file_type == "pdfs":
        hash_result = hash_pdf_file(file_bytes)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    elif file_type == "tables":
        hash_result = hash_table_file(file_bytes, filename)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    elif file_type == "pptx":
        hash_result = hash_pptx_file(file_bytes)
        if hash_result:
            hash_result = json.dumps(hash_result)
    
    if not hash_result:
        hash_result = compute_exact_hash(file_bytes)
    
    return hash_result


def compute_similarity(hash1, hash2):
    """Compute similarity between two hashes (0-100%)."""
    try:
        if hash1 == hash2:
            return 100.0
        
        is_image_hash1 = False
        is_image_hash2 = False
        
        try:
            if isinstance(hash1, str) and hash1.startswith('{'):
                h1_dict = json.loads(hash1)
                if 'phash' in h1_dict and 'color_hist' in h1_dict:
                    is_image_hash1 = True
        except:
            pass
        
        try:
            if isinstance(hash2, str) and hash2.startswith('{'):
                h2_dict = json.loads(hash2)
                if 'phash' in h2_dict and 'color_hist' in h2_dict:
                    is_image_hash2 = True
        except:
            pass
        
        if is_image_hash1 and is_image_hash2:
            h1 = json.loads(hash1)
            h2 = json.loads(hash2)
            
            phash_sim = 0
            ahash_sim = 0
            dhash_sim = 0
            
            try:
                bin_p1 = bin(int(h1['phash'], 16))[2:].zfill(64)
                bin_p2 = bin(int(h2['phash'], 16))[2:].zfill(64)
                hamming_p = sum(c1 != c2 for c1, c2 in zip(bin_p1, bin_p2))
                phash_sim = (1 - hamming_p / 64) * 100
            except:
                pass
            
            try:
                bin_a1 = bin(int(h1['ahash'], 16))[2:].zfill(64)
                bin_a2 = bin(int(h2['ahash'], 16))[2:].zfill(64)
                hamming_a = sum(c1 != c2 for c1, c2 in zip(bin_a1, bin_a2))
                ahash_sim = (1 - hamming_a / 64) * 100
            except:
                pass
            
            try:
                bin_d1 = bin(int(h1['dhash'], 16))[2:].zfill(64)
                bin_d2 = bin(int(h2['dhash'], 16))[2:].zfill(64)
                hamming_d = sum(c1 != c2 for c1, c2 in zip(bin_d1, bin_d2))
                dhash_sim = (1 - hamming_d / 64) * 100
            except:
                pass
            
            color_sim = 0
            try:
                hist1 = np.array(h1['color_hist'])
                hist2 = np.array(h2['color_hist'])
                
                norm1 = norm(hist1)
                norm2 = norm(hist2)
                
                if norm1 > 0 and norm2 > 0:
                    cosine_sim = np.dot(hist1, hist2) / (norm1 * norm2)
                    cosine_sim = np.clip(cosine_sim, 0, 1)
                    color_sim = cosine_sim * 100
            except:
                pass
            
            total_sim = (phash_sim * 0.4 + ahash_sim * 0.2 + dhash_sim * 0.2 + color_sim * 0.2)
            return float(total_sim)
        
        if isinstance(hash1, str) and hash1.startswith('[') and \
           isinstance(hash2, str) and hash2.startswith('['):
            vec1 = json.loads(hash1)
            vec2 = json.loads(hash2)
            
            vec1_filtered = np.array([v for v in vec1 if isinstance(v, (int, float))])
            vec2_filtered = np.array([v for v in vec2 if isinstance(v, (int, float))])
            
            if vec1_filtered.shape != vec2_filtered.shape or vec1_filtered.size == 0:
                return 0
            
            norm1 = norm(vec1_filtered)
            norm2 = norm(vec2_filtered)
            
            if norm1 == 0 or norm2 == 0:
                return 0
            
            cosine_sim = np.dot(vec1_filtered, vec2_filtered) / (norm1 * norm2)
            
            cosine_sim = np.clip(cosine_sim, 0, 1)
            return float(cosine_sim * 100)
        
        if (isinstance(hash1, str) and hash1.startswith('[')) or \
           (isinstance(hash2, str) and hash2.startswith('[')):
            return 0
        
        if is_image_hash1 or is_image_hash2:
            return 0
        
        try:
            bin_hash1 = bin(int(hash1, 16))[2:].zfill(128)
            bin_hash2 = bin(int(hash2, 16))[2:].zfill(128)
            hamming_dist = sum(c1 != c2 for c1, c2 in zip(bin_hash1, bin_hash2))
            similarity = (1 - hamming_dist / 128) * 100
            return float(similarity)
        except (ValueError, TypeError):
            pass
        
        if isinstance(hash1, str) and isinstance(hash2, str):
            if len(hash1) == len(hash2) and len(hash1) > 0:
                hamming = sum(c1 != c2 for c1, c2 in zip(hash1, hash2))
                similarity = (1 - hamming / len(hash1)) * 100
                return float(similarity)
        
        return 0.0
        
    except Exception as e:
        print(f"Similarity computation error: {e}")
        return 0.0


def detect_file_duplicates(file_records):
    """
    Detect duplicates using efficient pairwise comparison.
    
    file_records: list of dicts with keys:
        - id: file identifier
        - url: file URL (optional)
        - file_bytes: binary content
        - filename: original filename
    
    Returns: list of duplicate clusters
    """
    import time
    start_time = time.time()
    
    clusters = []
    seen_files = []
    
    files_by_type = defaultdict(list)
    for record in file_records:
        file_type = get_file_type(record.get('filename', ''))
        files_by_type[file_type].append(record)
    
    print(f"\n{'='*60}")
    print(f"Grouped files: {dict((k, len(v)) for k, v in files_by_type.items())}")
    print(f"{'='*60}")
    print(f"Processing {len(file_records)} total files...")
    print(f"Similarity threshold: {SIMILARITY_THRESHOLD}%")
    print(f"{'='*60}\n")
    
    for idx, record in enumerate(file_records):
        file_bytes = record.get('file_bytes')
        filename = record.get('filename', f'file_{idx}')
        
        if not file_bytes:
            print(f"[{idx+1}/{len(file_records)}] Skipping '{filename}' - no bytes")
            continue
        
        print(f"[{idx+1}/{len(file_records)}] Processing: {filename}")
        
        try:
            file_hash = compute_file_hash(file_bytes, filename)
        except Exception as e:
            print(f"  ⚠ Error computing hash: {e}")
            continue
        
        if not file_hash:
            print(f"  ⚠ Failed to compute hash")
            continue
        
        is_duplicate = False
        duplicate_cluster = None
        best_match = None
        best_similarity = 0
        
        for seen_idx, seen in enumerate(seen_files):
            try:
                similarity = compute_similarity(file_hash, seen["hash"])
                
                if similarity > best_similarity:
                    best_similarity = similarity
                    best_match = seen
                
                if similarity > 5:
                    print(f"    → vs '{seen['record'].get('filename')}': {similarity:.1f}%")
                
                if similarity >= SIMILARITY_THRESHOLD:  
                    print(f"  ✓ DUPLICATE DETECTED! Similarity: {similarity:.1f}% with '{seen['record'].get('filename')}'")
                    
                    for cluster in clusters:
                        if any(item['id'] == seen['record']['id'] for item in cluster):
                            duplicate_cluster = cluster
                            break
                    
                    if duplicate_cluster is None:
                        duplicate_cluster = [{
                            'id': seen['record']['id'],
                            'url': seen['record'].get('url'),
                            'filename': seen['record'].get('filename'),
                            'similarity_score': 1.0
                        }]
                        clusters.append(duplicate_cluster)
                    
                    duplicate_cluster.append({
                        'id': record['id'],
                        'url': record.get('url'),
                        'filename': record.get('filename'),
                        'similarity_score': round(similarity / 100, 2)
                    })
                    
                    is_duplicate = True
                    break
            except Exception as e:
                print(f"  ⚠ Error comparing with file {seen_idx}: {e}")
                continue
        
        if not is_duplicate:
            seen_files.append({
                'hash': file_hash,
                'record': record
            })
            if best_similarity >= SIMILARITY_THRESHOLD:
                print(f"  ⚠ Close match but at threshold: {best_similarity:.1f}%")
            elif best_similarity > 5:
                print(f"  ✓ Unique file (best match: {best_similarity:.1f}% - below {SIMILARITY_THRESHOLD}% threshold)")
            else:
                print(f"  ✓ Unique file (no similar matches found)")
    
    elapsed = time.time() - start_time
    
    print(f"\n{'='*60}")
    print(f"RESULTS:")
    print(f"  Similarity threshold: {SIMILARITY_THRESHOLD}%")
    print(f"  Total files processed: {len(file_records)}")
    print(f"  Unique files: {len(seen_files)}")
    print(f"  Duplicate clusters found: {len(clusters)}")
    print(f"  Total duplicates: {sum(len(c) - 1 for c in clusters)}")
    print(f"  Time taken: {elapsed:.2f}s")
    print(f"{'='*60}\n")
    
    return clusters